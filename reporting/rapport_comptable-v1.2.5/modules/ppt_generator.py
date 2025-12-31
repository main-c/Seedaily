#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Module 5: Génération PowerPoint

Génère une présentation PowerPoint complète basée sur le modèle existant
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
import openpyxl
from openpyxl.styles import numbers
import logging
from typing import Dict, Optional, Tuple
from datetime import datetime
import os
import subprocess
import tempfile
import time
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont

logger = logging.getLogger(__name__)


def apply_slide_template(slide, title: str, slide_number: int, total_slides: int,
                         periode: str = "", client: str = ""):
    """
    Applique un template professionnel uniforme à une slide

    Args:
        slide: La slide PowerPoint
        title: Titre de la slide
        slide_number: Numéro de la slide
        total_slides: Nombre total de slides
        periode: Période du rapport
        client: Nom du client
    """
    # Dimensions standard PowerPoint (10 x 7.5 inches)
    slide_width = Inches(10)
    slide_height = Inches(7.5)

    # 1. Header (bannière en haut)
    header = slide.shapes.add_shape(
        1,  # Rectangle
        0, 0,  # Position (x, y)
        slide_width,
        Inches(0.8)  # Hauteur du header
    )

    # Remplissage avec dégradé bleu
    fill = header.fill
    fill.gradient()
    fill.gradient_angle = 90.0
    fill.gradient_stops[0].color.rgb = RGBColor(54, 96, 146)  # Bleu foncé
    fill.gradient_stops[1].color.rgb = RGBColor(70, 120, 180)  # Bleu plus clair

    # Supprimer la bordure
    header.line.color.rgb = RGBColor(54, 96, 146)

    # Titre dans le header
    text_frame = header.text_frame
    text_frame.clear()
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_top = Inches(0.15)
    text_frame.vertical_anchor = 1  # Middle

    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = 'Calibri'

    # 2. Footer (pied de page)
    footer_height = Inches(0.4)
    footer = slide.shapes.add_shape(
        1,  # Rectangle
        0,
        slide_height - footer_height,
        slide_width,
        footer_height
    )

    # Remplissage gris clair
    fill = footer.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)

    # Bordure subtile en haut
    footer.line.color.rgb = RGBColor(200, 200, 200)
    footer.line.width = Pt(0.5)

    # Texte du footer: Période | Client | Page X/Y
    text_frame = footer.text_frame
    text_frame.clear()
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_right = Inches(0.3)
    text_frame.margin_top = Inches(0.08)

    p = text_frame.paragraphs[0]

    # Partie gauche: Période et Client
    left_text = f"{periode}  |  {client}" if periode and client else ""
    # Partie droite: Pagination
    right_text = f"Page {slide_number}/{total_slides}"

    # Combiner avec espaces pour alignement
    full_text = f"{left_text}{'':>{80}}{right_text}"
    p.text = full_text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.font.name = 'Calibri'

    # 3. Ligne de séparation décorative
    line = slide.shapes.add_connector(
        1,  # Straight connector
        Inches(0.5),
        Inches(0.85),
        Inches(9.5),
        Inches(0.85)
    )
    line.line.color.rgb = RGBColor(54, 96, 146)
    line.line.width = Pt(2)

    logger.debug(f"Template appliqué: {title} (Page {slide_number}/{total_slides})")


def add_styled_comment_box(slide, text: str, left: float, top: float,
                           width: float, height: float, icon: str = "📊"):
    """
    Ajoute un bloc de commentaires stylisé avec cadre et icône

    Args:
        slide: La slide PowerPoint
        text: Texte du commentaire
        left, top, width, height: Position et dimensions en inches
        icon: Icône à afficher (📊 analyse, 💡 recommandation, ⚠️ attention)
    """
    if not text or text.strip() == "":
        return  # Pas de commentaire à afficher

    # 1. Cadre principal (rectangle arrondi)
    from pptx.enum.shapes import MSO_SHAPE
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height)
    )

    # Fond bleu très clair
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(240, 248, 255)  # Alice Blue

    # Bordure gauche épaisse (accent)
    box.line.color.rgb = RGBColor(54, 96, 146)
    box.line.width = Pt(4)

    # 2. Icône dans un petit cercle en haut à gauche
    icon_size = 0.4
    icon_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left + 0.1),
        Inches(top + 0.1),
        Inches(icon_size),
        Inches(icon_size)
    )

    # Fond de l'icône
    icon_circle.fill.solid()
    icon_circle.fill.fore_color.rgb = RGBColor(54, 96, 146)
    icon_circle.line.fill.background()

    # Texte de l'icône
    icon_tf = icon_circle.text_frame
    icon_tf.text = icon
    icon_p = icon_tf.paragraphs[0]
    icon_p.font.size = Pt(18)
    icon_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    icon_tf.vertical_anchor = 1  # Middle

    # 3. Zone de texte pour le commentaire
    text_box = slide.shapes.add_textbox(
        Inches(left + 0.6),  # Décalage pour laisser place à l'icône
        Inches(top + 0.2),
        Inches(width - 0.8),
        Inches(height - 0.4)
    )

    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)

    # Parser et formater le texte
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]

        # Variable pour savoir si la couleur a déjà été définie
        color_set = False

        # Détecter les puces
        if line.strip().startswith('•') or line.strip().startswith('-'):
            p.text = line.strip()
            p.level = 0
            p.font.size = Pt(14)
            p.space_before = Pt(6)
        # Détecter les titres (tout en majuscules ou commence par nombre)
        elif line.strip().isupper() or (line.strip() and line.strip()[0].isdigit() and '.' in line[:3]):
            p.text = line.strip()
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(54, 96, 146)
            p.space_before = Pt(12)
            color_set = True  # Couleur déjà définie (bleu)
        else:
            p.text = line.strip()
            p.font.size = Pt(14)
            p.space_before = Pt(4)

        # Définir la police et la couleur par défaut si pas encore définie
        p.font.name = 'Calibri'
        if not color_set:
            p.font.color.rgb = RGBColor(50, 50, 50)

    logger.debug(f"Bloc commentaire stylisé ajouté avec icône {icon}")


def _excel_range_to_image(excel_path: str, sheet_name: str, cell_range: str, output_image: str) -> bool:
    """
    Exporte une plage Excel en image PNG
    - Sur Windows: utilise win32com (Excel COM automation)
    - Sur Linux/Mac: utilise LibreOffice Calc

    Args:
        excel_path: Chemin vers le fichier Excel
        sheet_name: Nom de la feuille
        cell_range: Plage de cellules (ex: "A1:F20")
        output_image: Chemin de sortie pour l'image PNG

    Returns:
        True si succès, False sinon
    """
    import platform

    # Déterminer la plateforme
    if platform.system() == 'Windows':
        return _excel_range_to_image_windows(excel_path, sheet_name, cell_range, output_image)
    else:
        return _excel_range_to_image_libreoffice(excel_path, sheet_name, cell_range, output_image)


def _excel_range_to_image_windows(excel_path: str, sheet_name: str, cell_range: str, output_image: str) -> bool:
    """
    Exporte une plage Excel en image PNG sur Windows en utilisant win32com
    """
    try:
        import win32com.client
        from PIL import ImageGrab

        # Créer une instance Excel
        excel = win32com.client.Dispatch("Excel.Application")
        excel.Visible = False
        excel.DisplayAlerts = False

        try:
            # Ouvrir le workbook
            wb = excel.Workbooks.Open(os.path.abspath(excel_path))
            ws = wb.Worksheets(sheet_name)

            # Sélectionner la plage
            range_obj = ws.Range(cell_range)

            # Copier la plage comme image dans le presse-papiers
            range_obj.CopyPicture(Format=2)  # 2 = xlBitmap

            # Capturer depuis le presse-papiers
            time.sleep(0.3)  # Petit délai pour s'assurer que le presse-papiers est prêt
            img = ImageGrab.grabclipboard()

            if img:
                # Sauvegarder l'image
                img.save(output_image, 'PNG')
                success = True
            else:
                logger.warning("Aucune image dans le presse-papiers")
                success = False

            # Fermer le workbook
            wb.Close(SaveChanges=False)

            return success

        finally:
            excel.Quit()

    except ImportError:
        logger.warning("win32com non disponible. Installez pywin32: pip install pywin32")
        return False
    except Exception as e:
        logger.error(f"Erreur export Excel Windows: {e}")
        return False


def _excel_range_to_image_libreoffice(excel_path: str, sheet_name: str, cell_range: str, output_image: str) -> bool:
    """
    Exporte une plage Excel en image PNG sur Linux/Mac
    Méthode: Crée un fichier Excel avec uniquement la plage, convertit en PDF, puis en PNG
    """
    import shutil

    temp_dir = None
    try:
        # Étape 1: Créer un nouveau fichier Excel avec seulement la plage voulue
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        ws = wb[sheet_name]

        # Parser la plage
        start_cell, end_cell = cell_range.split(':')
        min_row = ws[start_cell].row
        min_col = ws[start_cell].column
        max_row = ws[end_cell].row
        max_col = ws[end_cell].column

        # Créer un nouveau workbook avec seulement cette plage
        new_wb = openpyxl.Workbook()
        new_ws = new_wb.active
        new_ws.title = "Export"

        # Copier les données ET le formatage
        for i, row_idx in enumerate(range(min_row, max_row + 1), start=1):
            for j, col_idx in enumerate(range(min_col, max_col + 1), start=1):
                source_cell = ws.cell(row=row_idx, column=col_idx)
                target_cell = new_ws.cell(row=i, column=j)

                # Copier valeur
                target_cell.value = source_cell.value

                # Copier formatage
                if source_cell.font:
                    target_cell.font = source_cell.font.copy()
                if source_cell.fill:
                    target_cell.fill = source_cell.fill.copy()
                if source_cell.alignment:
                    target_cell.alignment = source_cell.alignment.copy()
                if source_cell.number_format:
                    target_cell.number_format = source_cell.number_format
                if source_cell.border:
                    target_cell.border = source_cell.border.copy()

        # Copier largeur des colonnes
        for j, col_idx in enumerate(range(min_col, max_col + 1), start=1):
            col_letter_src = openpyxl.utils.get_column_letter(col_idx)
            col_letter_dst = openpyxl.utils.get_column_letter(j)
            if ws.column_dimensions[col_letter_src].width:
                new_ws.column_dimensions[col_letter_dst].width = ws.column_dimensions[col_letter_src].width

        # Copier hauteur des lignes
        for i, row_idx in enumerate(range(min_row, max_row + 1), start=1):
            if ws.row_dimensions[row_idx].height:
                new_ws.row_dimensions[i].height = ws.row_dimensions[row_idx].height

        wb.close()

        # Sauvegarder le nouveau fichier
        temp_dir = tempfile.mkdtemp()
        temp_excel = os.path.join(temp_dir, 'range_export.xlsx')
        new_wb.save(temp_excel)
        new_wb.close()

        # Étape 2: Convertir en PDF avec LibreOffice
        result = subprocess.run([
            'libreoffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', temp_dir,
            temp_excel
        ], capture_output=True, timeout=30, text=True)

        temp_pdf = os.path.join(temp_dir, 'range_export.pdf')

        if result.returncode != 0 or not os.path.exists(temp_pdf):
            logger.debug(f"LibreOffice conversion PDF échouée: {result.stderr}")
            return False

        # Étape 3: Convertir le PDF en PNG
        from pdf2image import convert_from_path

        images = convert_from_path(
            temp_pdf,
            dpi=200,  # Haute résolution
            first_page=1,
            last_page=1,
            fmt='png'
        )

        if images:
            # Rogner les marges blanches
            img = images[0]

            # Optionnel: rogner les bords blancs
            from PIL import ImageChops
            bg = Image.new(img.mode, img.size, img.getpixel((0,0)))
            diff = ImageChops.difference(img, bg)
            bbox = diff.getbbox()
            if bbox:
                img = img.crop(bbox)

            # Sauvegarder
            img.save(output_image, 'PNG', quality=95, optimize=True)
            logger.info(f"✅ Tableau image: {sheet_name} [{cell_range}] exporté")
            return True
        else:
            return False

    except ImportError as e:
        logger.debug(f"Dépendance manquante: {e}")
        return False
    except subprocess.TimeoutExpired:
        logger.debug("Timeout LibreOffice")
        return False
    except Exception as e:
        logger.debug(f"Erreur export image: {e}")
        return False
    finally:
        # Nettoyer
        if temp_dir and os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir)
            except Exception:
                pass


def insert_excel_table_compact(slide, excel_path: str, sheet_name: str,
                               cell_range: str, left: float, top: float,
                               width: float, height: float, font_size: int = 8):
    """
    Insère un tableau Excel dans PowerPoint
    Méthode simple: copie les données et crée un tableau PowerPoint natif
    """
    try:
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        ws = wb[sheet_name]

        # Parser la plage
        start_cell, end_cell = cell_range.split(':')
        min_row = ws[start_cell].row
        min_col = ws[start_cell].column
        max_row = ws[end_cell].row
        max_col = ws[end_cell].column

        rows = max_row - min_row + 1
        cols = max_col - min_col + 1

        # Créer un tableau PowerPoint simple
        table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                             Inches(width), Inches(height))
        table = table_shape.table

        # Remplir le tableau
        for i, row_idx in enumerate(range(min_row, max_row + 1)):
            for j, col_idx in enumerate(range(min_col, max_col + 1)):
                cell = ws.cell(row=row_idx, column=col_idx)
                ppt_cell = table.cell(i, j)

                # Valeur
                value = cell.value
                if value is None or value == '':
                    ppt_cell.text = ""
                elif isinstance(value, (int, float)):
                    ppt_cell.text = f"{value:,.0f}" if abs(value) > 100 else f"{value:,.2f}"
                else:
                    ppt_cell.text = str(value)[:50]

                # Formatage de base
                tf = ppt_cell.text_frame
                tf.margin_top = Pt(2)
                tf.margin_bottom = Pt(2)
                tf.margin_left = Pt(4)
                tf.margin_right = Pt(4)
                tf.word_wrap = False

                para = tf.paragraphs[0]
                para.font.size = Pt(font_size)
                para.font.name = 'Calibri'

                # En-tête
                if i == 0:
                    ppt_cell.fill.solid()
                    ppt_cell.fill.fore_color.rgb = RGBColor(54, 96, 146)
                    para.font.color.rgb = RGBColor(255, 255, 255)
                    para.font.bold = True
                    para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

                # Ligne TOTAL (fond jaune)
                elif value and isinstance(value, str) and 'TOTAL' in value.upper():
                    ppt_cell.fill.solid()
                    ppt_cell.fill.fore_color.rgb = RGBColor(255, 242, 204)
                    para.font.bold = True

                # Lignes alternées
                elif i % 2 == 0:
                    ppt_cell.fill.solid()
                    ppt_cell.fill.fore_color.rgb = RGBColor(245, 248, 252)
                else:
                    ppt_cell.fill.solid()
                    ppt_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

                # Nombres négatifs en rouge
                if isinstance(value, (int, float)):
                    para.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
                    if value < 0:
                        para.font.color.rgb = RGBColor(192, 0, 0)
                        para.font.bold = True

        logger.info(f"✅ Tableau: {sheet_name} ({rows}×{cols})")
        wb.close()

    except Exception as e:
        logger.error(f"❌ Erreur tableau {sheet_name}: {e}")
        # Ajouter un texte d'erreur
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = textbox.text_frame
        tf.text = f"[Tableau: {sheet_name}]"
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)


def generate_powerpoint(excel_path: str, output_path: str, commentaires: Optional[Dict] = None,
                       template_path: Optional[str] = None):
    """Génère le rapport PowerPoint complet basé sur le modèle"""

    logger.info("=" * 80)
    logger.info("GÉNÉRATION DU RAPPORT POWERPOINT")
    logger.info("=" * 80)

    try:
        # Informations
        periode = commentaires.get('periode', 'Septembre 2025') if commentaires else 'Septembre 2025'
        cabinet = commentaires.get('cabinet', '2BN CONSULTING') if commentaires else '2BN CONSULTING'
        client = commentaires.get('client', 'BAMBOO IMMO') if commentaires else 'BAMBOO IMMO'

        # Créer présentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # NOMBRE TOTAL DE SLIDES (pour pagination) - sera ajusté à la fin
        total_slides = 16  # Estimation initiale

        # 1. Page de titre (pas de template)
        add_slide_1_titre(prs, client, periode, cabinet)

        # 2. Sommaire
        add_slide_2_sommaire(prs, periode, client)

        # 3. Objectif
        add_slide_3_objectif(prs, periode, client)

        # 4. Événements significatifs
        evenements = commentaires.get('bilan', {}).get('commentaire', '') if commentaires else ''
        add_slide_4_evenements(prs, evenements, periode, client)

        # 5-6. Situation financière (Bilan)
        bilan_comment = commentaires.get('bilan', {}).get('commentaire', '') if commentaires else ''
        add_slide_5_6_bilan(prs, bilan_comment, excel_path, periode, client)

        # 7-8. Activité (Compte de Résultat)
        cr_comment = commentaires.get('compte_resultat', {}).get('commentaire', '') if commentaires else ''
        add_slide_7_8_activite(prs, cr_comment, excel_path, periode, client)

        # 9-10. SIG
        sig_comment = commentaires.get('sig', {}).get('commentaire', '') if commentaires else ''
        add_slide_9_10_sig(prs, sig_comment, excel_path, periode, client)

        # 11. Situation mensuelle (Suivi Activité)
        suivi_comment = commentaires.get('suivi_activite', {}).get('commentaire', '') if commentaires else ''
        add_slide_11_mensuel(prs, suivi_comment, excel_path, periode, client)

        # 12. Décisions / Synthèse
        synthese_comment = commentaires.get('synthese', {}).get('commentaire', '') if commentaires else ''
        add_slide_12_decisions(prs, synthese_comment, periode, client)

        # 13-16. Annexes avec données de la Balance
        add_slides_annexes(prs, excel_path, periode, client)

        # Mettre à jour le nombre total de slides
        total_slides = len(prs.slides)
        logger.info(f"📊 Nombre total de slides: {total_slides}")

        # Sauvegarder
        prs.save(output_path)

        logger.info(f"✅ PowerPoint généré: {output_path}")
        logger.info("✨ Design amélioré: Templates uniformes, tableaux stylisés, commentaires enrichis")

    except Exception as e:
        logger.error(f"❌ Erreur: {e}", exc_info=True)
        raise


def insert_excel_table(slide, excel_path: str, sheet_name: str,
                       cell_range: str, left: float, top: float,
                       width: float, height: float, font_size: int = 9):
    """
    Insère un tableau Excel dans une slide PowerPoint

    Args:
        slide: La slide PowerPoint
        excel_path: Chemin vers le fichier Excel
        sheet_name: Nom de la feuille Excel
        cell_range: Plage de cellules (ex: "A1:F30")
        left, top, width, height: Position et dimensions en inches
        font_size: Taille de police (default: 9pt)
    """
    try:
        # Charger le fichier Excel
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        ws = wb[sheet_name]

        # Parser la plage de cellules
        start_cell, end_cell = cell_range.split(':')
        min_row = ws[start_cell].row
        min_col = ws[start_cell].column
        max_row = ws[end_cell].row
        max_col = ws[end_cell].column

        rows = max_row - min_row + 1
        cols = max_col - min_col + 1

        # Créer un tableau PowerPoint
        table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                             Inches(width), Inches(height))
        table = table_shape.table

        # Remplir le tableau avec les données Excel
        for i, row_idx in enumerate(range(min_row, max_row + 1)):
            for j, col_idx in enumerate(range(min_col, max_col + 1)):
                cell = ws.cell(row=row_idx, column=col_idx)
                ppt_cell = table.cell(i, j)

                # Récupérer la valeur
                value = cell.value
                if value is None:
                    ppt_cell.text = ""
                elif isinstance(value, (int, float)):
                    # Formater les nombres avec séparateurs de milliers
                    ppt_cell.text = f"{value:,.2f}" if value != 0 else "0.00"
                else:
                    ppt_cell.text = str(value)

                # Appliquer le formatage
                text_frame = ppt_cell.text_frame
                text_frame.paragraphs[0].font.size = Pt(font_size)
                text_frame.paragraphs[0].font.name = 'Calibri'
                text_frame.margin_top = Pt(2)
                text_frame.margin_bottom = Pt(2)
                text_frame.margin_left = Pt(3)
                text_frame.margin_right = Pt(3)

                # Couleur de fond pour les en-têtes (première ligne)
                if i == 0:
                    ppt_cell.fill.solid()
                    ppt_cell.fill.fore_color.rgb = RGBColor(54, 96, 146)
                    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    text_frame.paragraphs[0].font.bold = True

                # Alignement des nombres à droite
                if isinstance(value, (int, float)):
                    text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

        logger.info(f"✅ Tableau inséré: {sheet_name} ({rows}x{cols}) police={font_size}pt")

    except Exception as e:
        logger.error(f"❌ Erreur insertion tableau {sheet_name}: {e}")
        # Ajouter un message d'erreur dans la slide
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = textbox.text_frame
        tf.text = f"[Tableau {sheet_name} non disponible]"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.italic = True


def add_slide_1_titre(prs, client, periode, cabinet):
    """Diapo 1: Page de titre"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(54, 96, 146)
    
    add_text(slide, client, 1, 1.5, 8, 1, 40, True, RGBColor(255,255,255))
    add_text(slide, "Document confidentiel", 1, 2.7, 8, 0.5, 14, False, RGBColor(220,220,220))
    add_text(slide, f"Soumis à l'attention du service comptable", 1, 3.2, 8, 0.5, 14, False, RGBColor(220,220,220))
    add_text(slide, f"Période: {periode}", 1, 4.5, 8, 0.7, 24, True, RGBColor(255,255,255))
    add_text(slide, f"Par {cabinet}", 1, 6.5, 8, 0.5, 16, False, RGBColor(200,200,200))


def add_slide_2_sommaire(prs, periode: str = "", client: str = ""):
    """Diapo 2: Sommaire"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Appliquer le template
    total_slides = 16  # Estimation
    apply_slide_template(slide, "SOMMAIRE", 2, total_slides, periode, client)

    sommaire = """1- Objectif du rapport
2- Événements significatifs
3- Situation financière
4- Activité de la période
5- Soldes significatifs de gestion
6- Situation mensuelle
7- Décisions et recommandations
8- Annexes"""

    textbox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(4.5))
    tf = textbox.text_frame
    tf.text = sommaire
    for p in tf.paragraphs:
        p.font.size = Pt(20)
        p.space_after = Pt(14)
        p.font.name = 'Calibri'


def add_slide_3_objectif(prs, periode: str = "", client: str = ""):
    """Diapo 3: Objectif"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Appliquer le template avec header/footer
    apply_slide_template(slide, "1- Objectif du rapport", 3, 16, periode, client)

    # Nom du client dynamique
    client_name = client if client else "l'entreprise"

    objectif = f"""Le présent rapport a pour but :

• D'apprécier la situation financière de {client_name} ;
• De déterminer la profitabilité de {client_name} sur cette période ;
• Servir d'outil d'aide à la décision au Senior management."""

    # Ajouter la zone de texte en dessous du header
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    tf = textbox.text_frame
    tf.text = objectif
    tf.word_wrap = True

    for p in tf.paragraphs:
        p.font.size = Pt(18)
        p.font.name = 'Calibri'
        p.space_after = Pt(12)


def add_slide_4_evenements(prs, evenements, periode: str = "", client: str = ""):
    """Diapo 4: Événements significatifs"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    titre = "2- Événements significatifs"
    apply_slide_template(slide, titre, 4, 16, periode, client)

    text = evenements if evenements else "Aucun événement significatif à signaler pour cette période."
    add_styled_comment_box(slide, text, 1, 1.5, 8, 5, "⚠️")  # Icône attention pour événements


def add_slide_5_6_bilan(prs, commentaire, excel_path=None, periode: str = "", client: str = ""):
    """Diapos 5-6: Situation financière (Bilan)"""
    # Diapo 5: Tableau BILAN SYNTH - position selon rapport original
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Titre dynamique avec la période
    titre = f"3- Situation financière - {periode}" if periode else "3- Situation financière"
    apply_slide_template(slide, titre, 5, 16, periode, client)

    # Insérer tableau compact (ajuster position pour header)
    if excel_path and os.path.exists(excel_path):
        insert_excel_table_compact(slide, excel_path, "BILAN SYNTH", "A1:F30",
                                  left=0.1, top=1.2, width=9.8, height=5.8, font_size=9)
    else:
        add_text(slide, "BILAN SYNTHÉTIQUE", 0.5, 3.5, 9, 1, 28, True, RGBColor(80,80,80))

    # Diapo 6: Commentaires
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_template(slide, titre, 6, 16, periode, client)

    text = commentaire if commentaire else "Analyse du bilan à compléter."
    add_styled_comment_box(slide, text, 1, 1.5, 8, 5, "📊")


def add_slide_7_8_activite(prs, commentaire, excel_path=None, periode: str = "", client: str = ""):
    """Diapos 7-8: Activité (CR)"""
    # Diapo 7: Tableau COMPTE DE RÉSULTAT
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    titre = f"4- Activité {client} - {periode}" if client and periode else "4- Activité de la période"
    apply_slide_template(slide, titre, 7, 16, periode, client)

    if excel_path and os.path.exists(excel_path):
        insert_excel_table_compact(slide, excel_path, "CR SYNTH", "A1:F42",
                                  left=0.13, top=1.2, width=9.75, height=5.7, font_size=8)
    else:
        add_text(slide, "COMPTE DE RÉSULTAT", 0.5, 3.5, 9, 1, 28, True, RGBColor(80,80,80))

    # Diapo 8: Commentaires
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_template(slide, titre, 8, 16, periode, client)

    text = commentaire if commentaire else "Analyse du compte de résultat à compléter."
    add_styled_comment_box(slide, text, 1, 1.5, 8, 5, "📊")


def add_slide_9_10_sig(prs, commentaire, excel_path=None, periode: str = "", client: str = ""):
    """Diapos 9-10: SIG"""
    # Diapo 9: Tableau SIG
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    titre = "5- Soldes intermédiaires de gestion"
    apply_slide_template(slide, titre, 9, 16, periode, client)

    if excel_path and os.path.exists(excel_path):
        insert_excel_table_compact(slide, excel_path, "SIG", "A1:F44",
                                  left=0.16, top=1.1, width=9.7, height=5.9, font_size=7)
    else:
        add_text(slide, "INDICATEURS CLÉS", 0.5, 3.5, 9, 1, 28, True, RGBColor(80,80,80))

    # Diapo 10: Commentaires
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_template(slide, titre, 10, 16, periode, client)

    text = commentaire if commentaire else "Analyse des SIG à compléter."
    add_styled_comment_box(slide, text, 1, 1.5, 8, 5, "📊")


def add_slide_11_mensuel(prs, commentaire, excel_path=None, periode: str = "", client: str = ""):
    """Diapo 11: Situation mensuelle"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Extraire l'année de la période
    import re
    annee = re.search(r'\d{4}', periode) if periode else None
    annee_str = annee.group() if annee else "2025"

    titre = f"6- Situation mensuelle {client} {annee_str}" if client else f"6- Situation mensuelle {annee_str}"
    apply_slide_template(slide, titre, 11, 16, periode, client)

    if excel_path and os.path.exists(excel_path):
        insert_excel_table_compact(slide, excel_path, "SUIVI ACTIVITE", "A1:M66",
                                  left=0.05, top=1.1, width=9.9, height=5.9, font_size=6)
    else:
        text = commentaire if commentaire else "Suivi mensuel de l'activité à compléter."
        add_styled_comment_box(slide, text, 1, 1.5, 8, 5, "📊")


def add_slide_12_decisions(prs, synthese, periode: str = "", client: str = ""):
    """Diapo 12: Décisions"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    titre = "7- Décisions et Recommandations"
    apply_slide_template(slide, titre, 12, 16, periode, client)

    text = synthese if synthese else "Synthèse et recommandations à compléter."
    add_styled_comment_box(slide, text, 1, 1.5, 8, 5, "💡")  # Icône ampoule pour recommandations


def add_slides_annexes(prs, excel_path=None, periode: str = "", client: str = ""):
    """Diapos 13-16: Annexes simplifiées avec petits tableaux"""

    # Annexe 1-4: Tableaux compacts selon rapport original
    annexes = [
        ("Annexe 1 : Décaissements à justifier", ["47"], 0.21, 1.5, 9.63, 4.5, 13),
        ("Annexe 2 : Dettes Fournisseurs", ["401"], 0.18, 1.5, 9.64, 5.0, 14),
        ("Annexe 3 : Autres dettes, Préfinancement", ["42", "43", "44"], 0.76, 1.5, 8.55, 5.0, 15),
        ("Annexe 4 : Créditeurs divers", ["471", "472", "473"], 0.19, 1.5, 9.63, 4.5, 16)
    ]

    for titre, prefixes, left, top, width, height, slide_num in annexes:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_slide_template(slide, titre, slide_num, 16, periode, client)

        if excel_path and os.path.exists(excel_path):
            add_annexe_table_compact(slide, excel_path, prefixes, left, top, width, height)
        else:
            add_text(slide, "(Détails à compléter)", 0.5, 3.5, 9, 1, 18, False, RGBColor(120,120,120))


def add_annexe_table_compact(slide, excel_path: str, prefixes: list,
                             left: float, top: float, width: float, height: float):
    """Ajoute un tableau d'annexe compact filtré par préfixes de comptes"""
    try:
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        ws = wb["BG BI SEP"]

        # Filtrer les comptes
        filtered_rows = []
        for row_idx in range(4, ws.max_row + 1):
            compte = ws.cell(row_idx, 1).value
            if compte and any(str(compte).startswith(prefix) for prefix in prefixes):
                row_data = [ws.cell(row_idx, col_idx).value for col_idx in range(1, 7)]
                filtered_rows.append(row_data)

        if not filtered_rows:
            add_text(slide, "Aucun compte trouvé", left, top, width, height, 14, False, RGBColor(120,120,120))
            return

        # Créer petit tableau compact
        rows = len(filtered_rows) + 1
        table_shape = slide.shapes.add_table(rows, 6, Inches(left), Inches(top),
                                             Inches(width), Inches(height))
        table = table_shape.table

        # En-tête
        headers = ['N° Cpte', 'Libellé', 'Débit', 'Crédit', 'S. Déb', 'S. Créd']
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(54, 96, 146)
            tf = cell.text_frame
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.size = Pt(9)
            tf.margin_top = Pt(1)
            tf.margin_bottom = Pt(1)

        # Données
        for i, row_data in enumerate(filtered_rows, 1):
            for j, value in enumerate(row_data):
                cell = table.cell(i, j)
                tf = cell.text_frame
                tf.margin_top = Pt(1)
                tf.margin_bottom = Pt(1)

                if value is None or value == '':
                    cell.text = ""
                elif isinstance(value, (int, float)):
                    cell.text = f"{value:,.0f}"
                    tf.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
                else:
                    cell.text = str(value)[:30]

                tf.paragraphs[0].font.size = Pt(8)
                tf.paragraphs[0].font.name = 'Calibri'

        logger.info(f"✅ Annexe: {len(filtered_rows)} comptes")

    except Exception as e:
        logger.error(f"❌ Erreur annexe: {e}")
        add_text(slide, "[Erreur chargement]", left, top, width, height, 14, False, RGBColor(200, 50, 50))


def add_text(slide, text, left, top, width, height, size, bold, color):
    """Ajoute une zone de texte"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = textbox.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    generate_powerpoint("test.xlsx", "test.pptx")
